from fastapi import FastAPI, UploadFile, File, Query
from fastapi.responses import JSONResponse
from typing import List
import os
import shutil
import fitz  # PyMuPDF
import docx
import pptx
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

app = FastAPI()



app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Load embedding model
model = SentenceTransformer('all-MiniLM-L6-v2')

# Initialize FAISS index
embedding_dim = 384  # Dimension for MiniLM
index = faiss.IndexFlatL2(embedding_dim)

# In-memory store to map vectors to document texts
doc_texts = []

# Utility functions to extract text
def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# File upload endpoint
@app.post("/upload/")
async def upload_files(files: List[UploadFile] = File(...)):
    global doc_texts

    # Ensure 'temp' directory exists
    if not os.path.exists("temp"):
        os.makedirs("temp")

    for file in files:
        file_location = f"temp/{file.filename}"
        with open(file_location, "wb") as f:
            shutil.copyfileobj(file.file, f)

        ext = file.filename.split('.')[-1].lower()
        if ext == 'pdf':
            text = extract_text_from_pdf(file_location)
        elif ext == 'docx':
            text = extract_text_from_docx(file_location)
        elif ext == 'pptx':
            text = extract_text_from_pptx(file_location)
        elif ext == 'txt':
            with open(file_location, 'r', encoding='utf-8') as f:
                text = f.read()
        else:
            continue

        # Embed and add to FAISS
        embedding = model.encode([text])
        index.add(np.array(embedding, dtype='float32'))
        doc_texts.append(text)

        os.remove(file_location)

    return {"status": "Files processed and indexed."}

# Simple MCQ Generator from Text
def generate_mcq(text, num_mcqs):
    sentences = [s.strip() for s in text.replace('\n', ' ').split(". ") if len(s.split()) > 5]
    mcqs = []
    for i, sentence in enumerate(sentences[:num_mcqs]):
        question = sentence.strip(".?")
        choices = [question]
        # Generate 3 distractors (mocked here)
        for j in range(3):
            choices.append(f"Distractor {j+1} for '{question[:20]}...'")
        np.random.shuffle(choices)
        mcqs.append({
            "question": question,
            "choices": choices,
            "answer": question
        })
    return mcqs

# Query endpoint to get MCQs
@app.get("/mcqs/{num_mcqs}")
async def get_mcqs(num_mcqs: int):
    if not doc_texts:
        return JSONResponse(content={"error": "No documents uploaded."}, status_code=404)

    matched_text = doc_texts[-1]  # Use the last uploaded document
    mcqs = generate_mcq(matched_text, num_mcqs)
    return {"mcqs": mcqs}

if __name__ == "__main__":
    uvicorn.run(app, host="127.0.0.1", port=8000)
